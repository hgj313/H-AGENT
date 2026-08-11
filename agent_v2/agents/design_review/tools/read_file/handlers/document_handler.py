"""Document File Handler

Handler for document files (PDF, Word, Excel, etc.)
"""

from pathlib import Path

from .base_handler import BaseFileHandler, FileReadResult, HandlerCapability


class DocumentFileHandler(BaseFileHandler):
    """Handler for document files
    
    Supports: .pdf, .doc, .docx, .xls, .xlsx, .ppt, .pptx, etc.
    """
    
    name = "document"
    supported_extensions = {
        '.pdf', '.doc', '.docx',
        '.xls', '.xlsx', '.xlsm', '.xlsb',
        '.ppt', '.pptx', '.odp',
        '.rtf', '.odt', '.pages',
    }
    
    def _read_impl(self, file_path: str) -> FileReadResult:
        """Read document file
        
        Args:
            file_path: File path
            
        Returns:
            FileReadResult
        """
        path = Path(file_path)
        extension = path.suffix.lower()
        
        if extension == '.pdf':
            return self._read_pdf(file_path)
        elif extension in {'.doc', '.docx', '.rtf', '.odt', '.pages'}:
            return self._read_word(file_path)
        elif extension in {'.xls', '.xlsx', '.xlsm', '.xlsb'}:
            return self._read_excel(file_path)
        elif extension in {'.ppt', '.pptx', '.odp'}:
            return self._read_powerpoint(file_path)
        else:
            return FileReadResult(
                success=False,
                content="",
                error=f"不支持的文档格式: {extension}"
            )
    
    def _read_pdf(self, file_path: str) -> FileReadResult:
        """Read PDF file
        
        Args:
            file_path: File path
            
        Returns:
            FileReadResult
        """
        try:
            import PyPDF2
            
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = []
                
                for page in reader.pages:
                    text.append(page.extract_text())
                
                content = '\n\n'.join(text)
            
            return FileReadResult(
                success=True,
                content=content,
                metadata=self.get_metadata(file_path)
            )
        except ImportError:
            return self._read_pdf_with_pypdf(file_path)
        except Exception as e:
            return FileReadResult(
                success=False,
                content="",
                error=f"读取PDF失败: {str(e)}"
            )
    
    def _read_pdf_with_pypdf(self, file_path: str) -> FileReadResult:
        """Alternative PDF reading with pypdf"""
        try:
            import pypdf
            
            reader = pypdf.PdfReader(file_path)
            text = []
            
            for page in reader.pages:
                text.append(page.extract_text())
            
            content = '\n\n'.join(text)
            
            return FileReadResult(
                success=True,
                content=content,
                metadata=self.get_metadata(file_path)
            )
        except ImportError:
            return FileReadResult(
                success=False,
                content="",
                error="请安装 pypdf2 或 pypdf 库来读取 PDF 文件"
            )
        except Exception as e:
            return FileReadResult(
                success=False,
                content="",
                error=f"读取PDF失败: {str(e)}"
            )
    
    def _read_word(self, file_path: str) -> FileReadResult:
        """Read Word document
        
        Args:
            file_path: File path
            
        Returns:
            FileReadResult
        """
        try:
            from docx import Document
            
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs]
            content = '\n'.join(paragraphs)
            
            tables_content = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    tables_content.append('\t'.join(row_text))
            
            if tables_content:
                content += '\n\n[Tables]\n' + '\n'.join(tables_content)
            
            return FileReadResult(
                success=True,
                content=content,
                metadata=self.get_metadata(file_path)
            )
        except ImportError:
            return FileReadResult(
                success=False,
                content="",
                error="请安装 python-docx 库来读取 Word 文档"
            )
        except Exception as e:
            return FileReadResult(
                success=False,
                content="",
                error=f"读取Word文档失败: {str(e)}"
            )
    
    def _read_excel(self, file_path: str) -> FileReadResult:
        """Read Excel file
        
        Args:
            file_path: File path
            
        Returns:
            FileReadResult
        """
        try:
            import pandas as pd
            
            sheets = pd.read_excel(file_path, sheet_name=None)
            content_parts = []
            
            for sheet_name, df in sheets.items():
                content_parts.append(f"[{sheet_name}]\n")
                content_parts.append(df.to_string())
                content_parts.append("\n\n")
            
            content = ''.join(content_parts)
            
            return FileReadResult(
                success=True,
                content=content,
                metadata=self.get_metadata(file_path)
            )
        except ImportError:
            return FileReadResult(
                success=False,
                content="",
                error="请安装 pandas 和 openpyxl 库来读取 Excel 文件"
            )
        except Exception as e:
            return FileReadResult(
                success=False,
                content="",
                error=f"读取Excel文件失败: {str(e)}"
            )
    
    def _read_powerpoint(self, file_path: str) -> FileReadResult:
        """Read PowerPoint file
        
        Args:
            file_path: File path
            
        Returns:
            FileReadResult
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            content_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                content_parts.append(f"\n[Slide {slide_num}]\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content_parts.append(shape.text + "\n")
            
            content = ''.join(content_parts)
            
            return FileReadResult(
                success=True,
                content=content,
                metadata=self.get_metadata(file_path)
            )
        except ImportError:
            return FileReadResult(
                success=False,
                content="",
                error="请安装 python-pptx 库来读取 PowerPoint 文件"
            )
        except Exception as e:
            return FileReadResult(
                success=False,
                content="",
                error=f"读取PowerPoint文件失败: {str(e)}"
            )
    
    def _get_capabilities(self) -> set[HandlerCapability]:
        """Get handler capabilities"""
        return {
            HandlerCapability.READ_TEXT,
            HandlerCapability.SUPPORT_LARGE_FILE,
        }