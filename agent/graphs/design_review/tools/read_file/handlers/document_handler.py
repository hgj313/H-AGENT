"""文档文件处理器。

支持 PDF、DOC、DOCX、XLS、XLSX、PPT、PPTX 等文档格式。
提取文本内容用于模型分析。
"""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Any

from .base_handler import (
    BaseFileHandler,
    FileReadResult,
    HandlerCapability,
)

logger = logging.getLogger(__name__)


class DocumentFileHandler(BaseFileHandler):
    name: str = "DocumentFileHandler"
    supported_extensions: set[str] = {
        '.pdf',
        '.doc', '.docx',
        '.xls', '.xlsx',
        '.ppt', '.pptx',
        '.odt', '.ods', '.odp',
        '.rtf',
        '.epub', '.mobi',
    }

    def __init__(self, max_file_size: int | None = None) -> None:
        super().__init__(max_file_size)

    def get_capabilities(self) -> set[HandlerCapability]:
        return {HandlerCapability.RAW_TEXT, HandlerCapability.METADATA}

    def _do_read(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        ext = file_path.suffix.lower()

        handlers = {
            '.pdf': self._read_pdf,
            '.doc': self._read_doc,
            '.docx': self._read_docx,
            '.xls': self._read_xls,
            '.xlsx': self._read_xlsx,
            '.ppt': self._read_ppt,
            '.pptx': self._read_pptx,
            '.odt': self._read_odt,
            '.ods': self._read_ods,
            '.odp': self._read_odp,
            '.rtf': self._read_rtf,
            '.epub': self._read_epub,
            '.mobi': self._read_mobi,
        }

        handler = handlers.get(ext, self._read_generic)
        return handler(file_path, **kwargs)

    def _read_pdf(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        try:
            import pypdf

            reader = pypdf.PdfReader(file_path)
            metadata = {
                'page_count': len(reader.pages),
                'encrypted': reader.is_encrypted,
            }

            if reader.metadata:
                metadata['title'] = reader.metadata.get('/Title', '')
                metadata['author'] = reader.metadata.get('/Author', '')
                metadata['subject'] = reader.metadata.get('/Subject', '')
                metadata['creator'] = reader.metadata.get('/Creator', '')

            pages_text: list[str] = []
            start_page = kwargs.get('start_page', 0)
            end_page = kwargs.get('end_page', len(reader.pages))

            for i, page in enumerate(reader.pages[start_page:end_page]):
                try:
                    text = page.extract_text()
                    if text:
                        pages_text.append(f"[第 {start_page + i + 1} 页]\n{text}")
                except Exception as e:
                    self._logger.warning(f"提取第 {i + 1} 页文本失败: {e}")
                    pages_text.append(f"[第 {start_page + i + 1} 页]\n[无法提取文本]")

            full_text = "\n\n".join(pages_text)
            metadata['pages_extracted'] = len(pages_text)
            metadata['text_length'] = len(full_text)

            return self._create_success_result(
                content=full_text,
                file_path=file_path,
                capability=HandlerCapability.RAW_TEXT,
                metadata=metadata,
            )

        except ImportError:
            self._logger.warning("pypdf 未安装，尝试使用 PyPDF2")
            try:
                import PyPDF2

                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    metadata = {'page_count': len(reader.pages)}

                    pages_text: list[str] = []
                    for i, page in enumerate(reader.pages):
                        text = page.extract_text()
                        pages_text.append(f"[第 {i + 1} 页]\n{text}" if text else f"[第 {i + 1} 页]\n[无法提取文本]")

                    full_text = "\n\n".join(pages_text)
                    metadata['pages_extracted'] = len(pages_text)
                    metadata['text_length'] = len(full_text)

                    return self._create_success_result(
                        content=full_text,
                        file_path=file_path,
                        capability=HandlerCapability.RAW_TEXT,
                        metadata=metadata,
                    )
            except ImportError:
                return self._create_error_result("读取 PDF 需要安装 pypdf 或 PyPDF2", file_path)
        except Exception as e:
            self._logger.exception(f"读取 PDF 文件失败")
            return self._create_error_result(f"读取 PDF 失败: {str(e)}", file_path)

    def _read_docx(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        try:
            from docx import Document

            doc = Document(file_path)
            paragraphs: list[str] = []

            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)

            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(f"[表格行] {row_text}")

            full_text = "\n".join(paragraphs)
            metadata = {
                'paragraph_count': len(paragraphs),
                'table_count': len(doc.tables),
                'text_length': len(full_text),
            }

            return self._create_success_result(
                content=full_text,
                file_path=file_path,
                capability=HandlerCapability.RAW_TEXT,
                metadata=metadata,
            )

        except ImportError:
            return self._create_error_result("读取 DOCX 需要安装 python-docx", file_path)
        except Exception as e:
            self._logger.exception(f"读取 DOCX 文件失败")
            return self._create_error_result(f"读取 DOCX 失败: {str(e)}", file_path)

    def _read_doc(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        try:
            import antiword

            text = antiword.read(file_path)

            metadata = {
                'text_length': len(text),
            }

            return self._create_success_result(
                content=text,
                file_path=file_path,
                capability=HandlerCapability.RAW_TEXT,
                metadata=metadata,
            )

        except ImportError:
            return self._create_error_result("读取 DOC 文件需要安装 antiword (仅支持 Linux/Mac)", file_path)
        except Exception as e:
            self._logger.exception(f"读取 DOC 文件失败")
            return self._create_error_result(f"读取 DOC 失败: {str(e)}", file_path)

    def _read_xlsx(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, read_only=True, data_only=True)
            all_content: list[str] = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                all_content.append(f"[工作表: {sheet_name}]")

                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(v) if v is not None else '' for v in row]
                    row_text = " | ".join(row_values).strip()
                    if row_text:
                        all_content.append(row_text)

            full_text = "\n".join(all_content)
            metadata = {
                'sheet_count': len(wb.sheetnames),
                'sheets': wb.sheetnames,
                'text_length': len(full_text),
            }

            wb.close()
            return self._create_success_result(
                content=full_text,
                file_path=file_path,
                capability=HandlerCapability.RAW_TEXT,
                metadata=metadata,
            )

        except ImportError:
            return self._create_error_result("读取 XLSX 需要安装 openpyxl", file_path)
        except Exception as e:
            self._logger.exception(f"读取 XLSX 文件失败")
            return self._create_error_result(f"读取 XLSX 失败: {str(e)}", file_path)

    def _read_xls(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        try:
            import xlrd

            wb = xlrd.open_workbook(file_path)
            all_content: list[str] = []

            for sheet_idx in range(wb.nsheets):
                sheet = wb.sheet_by_index(sheet_idx)
                all_content.append(f"[工作表: {sheet.name}]")

                for row_idx in range(sheet.nrows):
                    row_values = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
                    row_text = " | ".join(row_values).strip()
                    if row_text:
                        all_content.append(row_text)

            full_text = "\n".join(all_content)
            metadata = {
                'sheet_count': wb.nsheets,
                'text_length': len(full_text),
            }

            return self._create_success_result(
                content=full_text,
                file_path=file_path,
                capability=HandlerCapability.RAW_TEXT,
                metadata=metadata,
            )

        except ImportError:
            return self._create_error_result("读取 XLS 需要安装 xlrd", file_path)
        except Exception as e:
            self._logger.exception(f"读取 XLS 文件失败")
            return self._create_error_result(f"读取 XLS 失败: {str(e)}", file_path)

    def _read_pptx(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            all_content: list[str] = []

            for slide_num, slide in enumerate(prs.slides, 1):
                all_content.append(f"\n{'=' * 40}\n[幻灯片 {slide_num}]\n{'=' * 40}")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        all_content.append(shape.text)

                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                            if row_text:
                                all_content.append(f"[表格] {row_text}")

            full_text = "\n".join(all_content)
            metadata = {
                'slide_count': len(prs.slides),
                'text_length': len(full_text),
            }

            return self._create_success_result(
                content=full_text,
                file_path=file_path,
                capability=HandlerCapability.RAW_TEXT,
                metadata=metadata,
            )

        except ImportError:
            return self._create_error_result("读取 PPTX 需要安装 python-pptx", file_path)
        except Exception as e:
            self._logger.exception(f"读取 PPTX 文件失败")
            return self._create_error_result(f"读取 PPTX 失败: {str(e)}", file_path)

    def _read_ppt(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        return self._create_error_result(
            "读取旧版 PPT 格式需要先转换为 PPTX 格式",
            file_path
        )

    def _read_odt(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        return self._read_odf_generic(file_path, 'odt', **kwargs)

    def _read_ods(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        return self._read_odf_generic(file_path, 'ods', **kwargs)

    def _read_odp(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        return self._read_odf_generic(file_path, 'odp', **kwargs)

    def _read_odf_generic(self, file_path: Path, fmt: str, **kwargs: Any) -> FileReadResult:
        try:
            import odf

            all_content: list[str] = []

            with open(file_path, 'rb') as f:
                from odf.opendocument import load
                from odf import text, teletype

                doc = load(file_path)

                if fmt == 'odt':
                    paragraphs = doc.getElementsByType(text.P)
                    for p in paragraphs:
                        para_text = teletype.extractText(p)
                        if para_text.strip():
                            all_content.append(para_text)

                elif fmt in ('ods', 'odp'):
                    from odf import table, dcuitls

                    if fmt == 'ods':
                        tables = doc.getElementsByType(table.Table)
                        for tbl in tables:
                            rows = tbl.getElementsByType(table.Row)
                            for row in rows:
                                cells = row.getElementsByType(table.TableCell)
                                row_text = " | ".join(teletype.extractText(c) for c in cells if teletype.extractText(c).strip())
                                if row_text:
                                    all_content.append(row_text)

                    elif fmt == 'odp':
                        from odf.presentation import Page, Notes
                        slides = doc.getElementsByType(Page)
                        for idx, slide in enumerate(slides, 1):
                            all_content.append(f"[幻灯片 {idx}]")
                            allParagraphs = slide.getElementsByType(text.P)
                            for p in allParagraphs:
                                para_text = teletype.extractText(p)
                                if para_text.strip():
                                    all_content.append(para_text)

            full_text = "\n".join(all_content)
            metadata = {
                'format': f'odf-{fmt}',
                'text_length': len(full_text),
            }

            return self._create_success_result(
                content=full_text,
                file_path=file_path,
                capability=HandlerCapability.RAW_TEXT,
                metadata=metadata,
            )

        except ImportError:
            return self._create_error_result(f"读取 ODF 格式需要安装 odfpy", file_path)
        except Exception as e:
            self._logger.exception(f"读取 ODF 文件失败")
            return self._create_error_result(f"读取 ODF 文件失败: {str(e)}", file_path)

    def _read_rtf(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        try:
            import striprtf

            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()

            text = striprtf.rtf_to_text(rtf_content)
            metadata = {
                'text_length': len(text),
            }

            return self._create_success_result(
                content=text,
                file_path=file_path,
                capability=HandlerCapability.RAW_TEXT,
                metadata=metadata,
            )

        except ImportError:
            return self._create_error_result("读取 RTF 需要安装 striprtf", file_path)
        except Exception as e:
            self._logger.exception(f"读取 RTF 文件失败")
            return self._create_error_result(f"读取 RTF 失败: {str(e)}", file_path)

    def _read_epub(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        try:
            import epublib

            all_content: list[str] = []
            book = epublib.Epub.open(str(file_path))

            for chapter in book.get_chapters():
                if chapter.content:
                    all_content.append(f"[{chapter.title}]\n{chapter.content}")

            full_text = "\n\n".join(all_content)
            metadata = {
                'title': getattr(book, 'title', ''),
                'author': getattr(book, 'author', ''),
                'chapter_count': len(all_content),
                'text_length': len(full_text),
            }

            return self._create_success_result(
                content=full_text,
                file_path=file_path,
                capability=HandlerCapability.RAW_TEXT,
                metadata=metadata,
            )

        except ImportError:
            return self._create_error_result("读取 EPUB 需要安装 epublib", file_path)
        except Exception as e:
            self._logger.exception(f"读取 EPUB 文件失败")
            return self._create_error_result(f"读取 EPUB 失败: {str(e)}", file_path)

    def _read_mobi(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        try:
            import mobi

            all_content: list[str] = []
            result = mobi.extract(str(file_path))

            if result and result.get('html'):
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(result['html'], 'html.parser')
                text = soup.get_text(separator='\n', strip=True)
                all_content.append(text)

            full_text = "\n\n".join(all_content)
            metadata = {
                'text_length': len(full_text),
            }

            return self._create_success_result(
                content=full_text,
                file_path=file_path,
                capability=HandlerCapability.RAW_TEXT,
                metadata=metadata,
            )

        except ImportError:
            return self._create_error_result("读取 MOBI 需要安装 mobi", file_path)
        except Exception as e:
            self._logger.exception(f"读取 MOBI 文件失败")
            return self._create_error_result(f"读取 MOBI 失败: {str(e)}", file_path)

    def _read_generic(self, file_path: Path, **kwargs: Any) -> FileReadResult:
        return self._create_error_result(
            f"不支持的文件格式: {file_path.suffix}",
            file_path
        )